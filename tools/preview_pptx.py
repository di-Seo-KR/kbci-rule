# -*- coding: utf-8 -*-
"""생성된 PPTX를 PNG 미리보기로 렌더링(QA용 경량 렌더러).

LibreOffice 없이 python-pptx 의 도형/텍스트/그림을 직접 읽어 Pillow 로 그립니다.
완벽한 재현이 아니라 레이아웃 점검(겹침/넘침/배치)용입니다.

사용법: python tools/preview_pptx.py <PPTX> [--outdir DIR] [--dpi 96]
"""
import argparse
import io
import os

from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_FILL

FONT_PATH = "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"  # 한글 포함 CJK
_FONT_CACHE = {}


def font(px):
    px = max(8, int(px))
    if px not in _FONT_CACHE:
        try:
            _FONT_CACHE[px] = ImageFont.truetype(FONT_PATH, px)
        except Exception:
            _FONT_CACHE[px] = ImageFont.load_default()
    return _FONT_CACHE[px]


def rgb(c):
    return (c[0], c[1], c[2])


def safe_fill(shape):
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            return rgb(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def safe_line(shape):
    try:
        if shape.line.fill.type == MSO_FILL.SOLID:
            w = shape.line.width.pt if shape.line.width else 1
            return rgb(shape.line.color.rgb), max(1, int(w))
    except Exception:
        pass
    return None, 0


def run_style(run, dpi):
    sz = run.font.size.pt if run.font.size else 12
    try:
        col = rgb(run.font.color.rgb)
    except Exception:
        col = (51, 51, 51)
    return {"px": sz * dpi / 72.0, "bold": bool(run.font.bold), "color": col}


def draw_text_frame(draw, tf, box, dpi):
    bx, by, bw, bh = box
    # 문단 -> 줄 단위로 래핑
    para_lines = []  # [(line_chars[(ch,style)], line_w, line_h)]
    for para in tf.paragraphs:
        align = para.alignment
        chars = []
        for r in para.runs:
            st = run_style(r, dpi)
            for ch in r.text:
                chars.append((ch, st))
        if not chars:
            para_lines.append(([], 0, font(12 * dpi / 72.0).size * 1.3, align))
            continue
        line = []
        lw = 0
        for ch, st in chars:
            f = font(st["px"])
            cw = f.getlength(ch)
            if lw + cw > bw and line:
                lh = max(font(s["px"]).size for _, s in line) * 1.32
                para_lines.append((line, lw, lh, align))
                line, lw = [], 0
            line.append((ch, st))
            lw += cw
        if line:
            lh = max(font(s["px"]).size for _, s in line) * 1.32
            para_lines.append((line, lw, lh, align))

    total_h = sum(l[2] for l in para_lines)
    anchor = tf.vertical_anchor
    if anchor == MSO_ANCHOR.MIDDLE:
        y = by + max(0, (bh - total_h) / 2)
    elif anchor == MSO_ANCHOR.BOTTOM:
        y = by + max(0, bh - total_h)
    else:
        y = by
    for line, lw, lh, align in para_lines:
        if align == PP_ALIGN.CENTER:
            x = bx + max(0, (bw - lw) / 2)
        elif align == PP_ALIGN.RIGHT:
            x = bx + max(0, bw - lw)
        else:
            x = bx
        for ch, st in line:
            f = font(st["px"])
            draw.text((x, y), ch, font=f, fill=st["color"])
            if st["bold"]:
                draw.text((x + 1, y), ch, font=f, fill=st["color"])
            x += f.getlength(ch)
        y += lh


def render_slide(slide, w_emu, h_emu, dpi):
    scale = dpi / 914400.0
    W, H = int(w_emu * scale), int(h_emu * scale)
    img = Image.new("RGB", (W, H), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    def box(sh):
        return (sh.left * scale, sh.top * scale, sh.width * scale, sh.height * scale)

    for sh in slide.shapes:
        x, y, bw, bh = box(sh)
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                pic = pic.resize((max(1, int(bw)), max(1, int(bh))))
                img.paste(pic, (int(x), int(y)))
            except Exception:
                pass
            continue
        # 도형(사각/원) 채움/선
        fill = safe_fill(sh)
        lcol, lw = safe_line(sh)
        is_oval = (st == MSO_SHAPE_TYPE.AUTO_SHAPE and
                   getattr(sh, "auto_shape_type", None) == MSO_SHAPE.OVAL)
        if fill or lcol:
            xy = [int(x), int(y), int(x + bw), int(y + bh)]
            if is_oval:
                draw.ellipse(xy, fill=fill, outline=lcol, width=lw or 1)
            else:
                draw.rectangle(xy, fill=fill, outline=lcol, width=lw or 1)
        # 텍스트
        if sh.has_text_frame and sh.text_frame.text.strip():
            pad = 2
            draw_text_frame(draw, sh.text_frame, (x + pad, y + pad, bw - 2 * pad, bh - 2 * pad), dpi)
    return img


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--outdir", default="/tmp/preview")
    ap.add_argument("--dpi", type=float, default=96)
    args = ap.parse_args()
    os.makedirs(args.outdir, exist_ok=True)
    prs = Presentation(args.pptx)
    w, h = prs.slide_width, prs.slide_height
    paths = []
    for i, slide in enumerate(prs.slides):
        im = render_slide(slide, w, h, args.dpi)
        p = os.path.join(args.outdir, f"slide_{i+1:02d}.png")
        im.save(p)
        paths.append(p)
    print("rendered:", len(paths))
    for p in paths:
        print(p)


if __name__ == "__main__":
    main()
