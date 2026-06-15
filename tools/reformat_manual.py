# -*- coding: utf-8 -*-
"""StepHow export PPTX를 KB 사규관리시스템 매뉴얼 v1.0 양식으로 변환하는 범용 도구.

입력: StepHow에서 내보낸 PPTX(표지 + 섹션 구분 + n/N 단계 슬라이드 구조)
출력: v1.0 표준 양식 덱(실제 KB 로고, 헤더 위계, 좌측 확대 이미지 + 빨간 번호 마커 +
      빨간 외곽 하이라이트, 우측 설명, 바닥글 로고 + 자동 페이지 번호)

사용법: python tools/reformat_manual.py SOURCE.pptx [--out OUT.pptx]
"""
import argparse
import os
import re
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from PIL import Image

from build_manual_pptx import (
    KB_YELLOW, KB_YELLOW_DK, GREY_DK, GREY, GREY_LT, GREY_BORDER, WHITE,
    SLIDE_W_IN, SLIDE_H_IN, EMU_PER_IN,
    add_text, rect, number_chip,
)
from build_stephow_pptx import marker, desc_box, _slide_number_field, MARK_RED
from pptx.dml.color import RGBColor

EMU = 914400
HILITE_RED = RGBColor(0xC0, 0x00, 0x00)
LOGO = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "assets", "kb_logo.png"))


# ====================================================== 추출
def _txt(sh):
    t = sh.text_frame.text if sh.has_text_frame else ""
    return t.replace('\xa0', ' ').strip()


def _is_page(t):
    return bool(re.match(r'^\d+\s*/\s*\d+$', t))


def _is_header(t):
    return '매뉴얼' in t and '_v' in t


def _is_section_lbl(t):
    return t.startswith('◦')


def _title_num(t):
    m = re.match(r'^(\d+)\.\s*(.+)$', t)
    return (int(m.group(1)), m.group(2).strip()) if m else None


def _color(sh, which):
    try:
        f = sh.fill if which == 'fill' else sh.line.fill
        if f.type == MSO_FILL.SOLID:
            return str((sh.fill if which == 'fill' else sh.line.color).rgb)
    except Exception:
        return None
    return None


def _is_hilite_color(hexs):
    if not hexs:
        return False
    if hexs in ('C00000', 'D93726', 'FF6F35'):
        return True
    try:
        r, g, b = int(hexs[0:2], 16), int(hexs[2:4], 16), int(hexs[4:6], 16)
    except Exception:
        return False
    return r > 180 and 30 <= g <= 175 and b < 130  # 주황/빨강 계열


def classify_desc(text):
    items = []
    for line in text.split('\n'):
        s = line.strip()
        if not s:
            continue
        if s[0] in '①②③④⑤⑥⑦⑧⑨⑩' or re.match(r'^\d+[\).]', s):
            items.append((s, 'item'))
        elif s[0] in '※-•' or s.startswith('비고'):
            items.append((s, 'sub'))
        else:
            items.append((s, 'lead'))
    return items


def extract(src):
    p = Presentation(src)
    slides = list(p.slides)
    meta = {'title': '매뉴얼', 'version': 'v0.1', 'dept': '', 'authors': '',
            'written': '', 'revised': '', 'contact': '', 'steps': 0}
    # 제목/버전: 러닝 헤더
    for s in slides:
        done = False
        for sh in s.shapes:
            m = re.match(r'^(.*)_v([\d.]+)$', _txt(sh))
            if m:
                meta['title'] = m.group(1).strip()
                meta['version'] = 'v' + m.group(2)
                done = True
                break
        if done:
            break
    # 표지 메타
    for sh in slides[0].shapes:
        for l in _txt(sh).split('\n'):
            l = l.strip()
            if l.startswith('작성 부서'):
                meta['dept'] = l.split(':', 1)[1].strip()
            elif l.startswith('담당자'):
                meta['authors'] = l.split(':', 1)[1].strip()
            elif l.startswith('작성') and '일자' in l:
                meta['written'] = l.split(':', 1)[1].strip()
            elif l.startswith('수정') and '일자' in l:
                meta['revised'] = l.split(':', 1)[1].strip()
            elif l.startswith('*') or ('문의' in l and ':' not in l[:3]):
                meta['contact'] = l.lstrip('*').strip()
            elif re.match(r'^\d+\s*단계$', l):
                meta['steps'] = int(re.match(r'^(\d+)', l).group(1))

    items = []
    for s in slides[1:]:
        texts = [_txt(sh) for sh in s.shapes if _txt(sh)]
        # 섹션 구분?
        is_sec = any(re.match(r'^섹션\s*\d+\s*/\s*\d+$', t) for t in texts)
        roman_t = None
        for t in texts:
            m = re.match(r'^([IVX]+)\.\s*(.+)$', t)
            if m:
                roman_t = (m.group(1), m.group(2).strip())
        if is_sec and roman_t:
            items.append(('section', roman_t[0], roman_t[1]))
            continue
        # 스텝
        pagenum = None
        for t in texts:
            mp = re.match(r'^(\d+)\s*/\s*\d+$', t)
            if mp:
                pagenum = int(mp.group(1))
        if pagenum is None:
            continue
        title = ''
        for t in texts:
            tn = _title_num(t)
            if tn and tn[0] == pagenum:
                title = tn[1]
        cands = [t for t in texts if not _is_page(t) and not _is_header(t)
                 and not _is_section_lbl(t)
                 and not (_title_num(t) and _title_num(t)[0] == pagenum)]
        desc = max(cands, key=len) if cands else ''
        # 이미지/마커/하이라이트
        pics = [sh for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        img, markers, hilites = None, [], []
        if pics:
            pic = max(pics, key=lambda sh: sh.width * sh.height)
            img = pic.image.blob
            px, py, pw, ph = pic.left, pic.top, pic.width, pic.height
            for sh in s.shapes:
                if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                    continue
                t = _txt(sh)
                fx = (sh.left + sh.width / 2 - px) / pw
                fy = (sh.top + sh.height / 2 - py) / ph
                if t.isdigit() and 0 <= fx <= 1 and 0 <= fy <= 1:
                    markers.append((int(t), round(fx, 4), round(fy, 4)))
                elif not t and sh.width > 0.3 * EMU and sh.height > 0.2 * EMU:
                    if _is_hilite_color(_color(sh, 'fill')) or _is_hilite_color(_color(sh, 'line')):
                        bx, by = (sh.left - px) / pw, (sh.top - py) / ph
                        bw, bh = sh.width / pw, sh.height / ph
                        if 0 <= bx <= 1 and 0 <= by <= 1 and 0.02 < bw * bh < 0.9:
                            hilites.append((bx, by, bw, bh, sh.width * sh.height))
            markers.sort()
            if hilites:
                hilites = [max(hilites, key=lambda h: h[4])[:4]]
        items.append(('step', {'n': pagenum, 'title': title, 'desc': desc,
                               'img': img, 'markers': markers, 'hilites': hilites}))
    return meta, items


# ====================================================== 렌더(v1.0 양식)
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def footer(slide):
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(0.31), Inches(7.0), Inches(1.91), Inches(0.31))
    tb = slide.shapes.add_textbox(Inches(SLIDE_W_IN - 1.35), Inches(6.98), Inches(0.9), Inches(0.32))
    tb.text_frame.word_wrap = False
    pp = tb.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    _slide_number_field(pp, 12, GREY)


def cover_slide(prs, meta):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(0.31), Inches(0.36), Inches(2.41), Inches(0.4))
    rect(s, 0, 2.0, SLIDE_W_IN, 1.95, fill=KB_YELLOW)
    rect(s, 0, 2.0, SLIDE_W_IN, 0.08, fill=KB_YELLOW_DK)
    add_text(s, 0.9, 2.32, 11.5, 0.9, [[(meta['title'], dict(size=36, bold=True, color=WHITE))]])
    add_text(s, 0.9, 3.34, 11.5, 0.5,
             [[(f"Ver {meta['version'].lstrip('vV')}", dict(size=20, bold=True, color=WHITE))]])
    bx, by, bw, bh = 0.68, 4.94, 7.6, 2.03
    rect(s, bx, by, bw, bh, fill=GREY_LT, line=GREY_BORDER, line_w=1)
    info = []
    if meta['dept']:
        info.append([("작성 부서 : ", dict(size=13, bold=True, color=GREY_DK)), (meta['dept'], dict(size=13, color=GREY_DK))])
    if meta['authors']:
        info.append([("담당자 : ", dict(size=13, bold=True, color=GREY_DK)), (meta['authors'], dict(size=13, color=GREY_DK))])
    wr = []
    if meta['written']:
        wr += [("작성 일자 : ", dict(size=13, bold=True, color=GREY_DK)), (meta['written'], dict(size=13, color=GREY_DK))]
    if meta['revised']:
        wr += [("    수정 일자 : ", dict(size=13, bold=True, color=GREY_DK)), (meta['revised'], dict(size=13, color=GREY_DK))]
    if wr:
        info.append(wr)
    if meta['contact']:
        info.append([(f"※ {meta['contact']}", dict(size=11, color=GREY))])
    add_text(s, bx + 0.25, by + 0.22, bw - 0.5, bh - 0.4, info)
    return s


def divider_slide(prs, roman, title):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    rect(s, 0, 2.75, SLIDE_W_IN, 2.0, fill=KB_YELLOW)
    rect(s, 0, 2.75, SLIDE_W_IN, 0.07, fill=KB_YELLOW_DK)
    add_text(s, 0.9, 3.05, 2.2, 1.4, [[(roman, dict(size=54, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 2.6, 3.05, 9.5, 1.4, [[(title, dict(size=34, bold=True, color=GREY_DK))]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def _header(s, n, title, section):
    rect(s, 0, 0, SLIDE_W_IN, 1.0, fill=KB_YELLOW)
    rect(s, 0, 1.0, SLIDE_W_IN, 0.05, fill=KB_YELLOW_DK)
    number_chip(s, 0.5, 0.27, 0.46, n, fill=WHITE, txt=KB_YELLOW_DK, size=13)
    add_text(s, 1.12, 0.1, 8.3, 0.8, [[(title, dict(size=22, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
    if section:
        add_text(s, 9.5, 0.1, 3.45, 0.8, [[(section, dict(size=13, bold=True, color=WHITE))]],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def step_slide(prs, st, section, tmp):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    n, title = st['n'], (st['title'] or '')
    items = classify_desc(st['desc'])
    _header(s, n, title, section)
    # 좌측 이미지 영역(틀 고정) — 이미지 없으면 자리표시자
    IMGX, IMGY, IMGW, IMGH = 0.1, 1.45, 8.4, 5.5
    if st['img']:
        ip = os.path.join(tmp, f"s{n}.png")
        open(ip, 'wb').write(st['img'])
        iw, ih = Image.open(ip).size
        scale = min(IMGW / (iw / 96.0), IMGH / (ih / 96.0))
        dw, dh = (iw / 96.0) * scale, (ih / 96.0) * scale
        dx, dy = IMGX + (IMGW - dw) / 2, IMGY + (IMGH - dh) / 2
        s.shapes.add_picture(ip, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
        rect(s, dx, dy, dw, dh, fill=None, line=GREY_BORDER, line_w=1)
        for bx, by, bw, bh in st['hilites']:
            rect(s, dx + bx * dw, dy + by * dh, bw * dw, bh * dh, fill=None, line=HILITE_RED, line_w=2.25)
        for mn, fx, fy in st['markers']:
            marker(s, dx + fx * dw, dy + fy * dh, 0.23, mn)
    else:
        rect(s, IMGX, IMGY, IMGW, IMGH, fill=GREY_LT, line=GREY_BORDER, line_w=1)
        add_text(s, IMGX, IMGY + IMGH / 2 - 0.35, IMGW, 0.7,
                 [[("[ 화면 캡처 영역 ]", dict(size=14, bold=True, color=GREY))],
                  [("캡처 추가 예정", dict(size=10, color=GREY))]],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    # 우측 '사용 방법' 패널(항상 유지 — 내용은 직접 작성)
    rect(s, 8.62, 1.5, 0.014, 5.4, fill=GREY_BORDER)
    RX, RW = 8.85, 4.3
    add_text(s, RX, 1.5, RW, 0.35, [[("사용 방법", dict(size=13, bold=True, color=KB_YELLOW_DK))]])
    rect(s, RX, 1.87, 0.55, 0.045, fill=KB_YELLOW)
    if items:
        desc_box(s, RX, 2.05, RW, 4.7, items)
    footer(s)
    return s


def build(src, out):
    meta, items = extract(src)
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * EMU_PER_IN))
    prs.slide_height = Emu(int(SLIDE_H_IN * EMU_PER_IN))
    cover_slide(prs, meta)
    cur = ('', '')
    n_steps = n_img = n_sec = 0
    with tempfile.TemporaryDirectory() as tmp:
        for it in items:
            if it[0] == 'section':
                cur = (it[1], it[2])
                divider_slide(prs, it[1], it[2])
                n_sec += 1
            else:
                section = f"{cur[0]}. {cur[1]}" if cur[0] else ''
                step_slide(prs, it[1], section, tmp)
                n_steps += 1
                n_img += 1 if it[1]['img'] else 0
        prs.save(out)
    print(f"생성: {out} | 표지1 + 섹션{n_sec} + 단계{n_steps}(이미지 {n_img}) = {1 + n_sec + n_steps}장")
    print(f"메타: {meta['title']} {meta['version']} / {meta['dept']} / {meta['authors']}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("src")
    ap.add_argument("--out", default="사용자매뉴얼/생성본/변환본.pptx")
    args = ap.parse_args()
    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    build(args.src, args.out)
