# -*- coding: utf-8 -*-
"""StepHow(가이드플로) 매뉴얼 HTML 내용을 기반으로 PPTX 생성.

- 텍스트(15단계·4섹션·표지)는 코드에 내장(원본 HTML에서 추출).
- 이미지는 guideflo 호스트 차단으로 직접 수집 불가 → captures/stephow/stepNN.png 가 있으면 삽입,
  없으면 자리표시자. (StepHow 캡처는 이미 주석이 입혀져 있어 별도 마커는 추가하지 않음)

사용법: python tools/build_stephow_pptx.py [--captures captures/stephow] [--out OUT.pptx]
"""
import argparse
import os
import tempfile
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# 검증된 테마/헬퍼 재사용
from build_manual_pptx import (
    KB_YELLOW, KB_YELLOW_DK, GREY_DK, GREY, GREY_LT, GREY_BORDER, WHITE,
    SLIDE_W_IN, SLIDE_H_IN, EMU_PER_IN, KR_FONT,
    style_run, add_text, rect, number_chip, build_image,
)
from stephow_markers import MARKERS

MARK_RED = RGBColor(0xD9, 0x37, 0x26)  # 번호 마커 색(고대비)

# ------------------------------------------------------------ 문서 메타/콘텐츠
DOC = {
    "title": "사규관리시스템 사용자 매뉴얼",
    "version": "v0.1",
    "org": "KB신용정보",
    "dept": "경영전략부",
    "authors": "박현민 과장, 서동인 계장",
    "written": "2026-06-15",
    "revised": "2026-06-15",
    "contact": "이용 관련 문의사항은 박현민 과장 또는 서동인 계장에게 문의 바랍니다.",
    "steps": 15,
}

PARTS = [
    ("I", "사규 검색", [1, 2, 3, 4, 5]),
    ("II", "사규 본문 조회", [6, 7, 8]),
    ("III", "사규 캘린더", [9, 10, 11]),
    ("IV", "법령 검색", [12, 13, 14, 15]),
]

# kind: lead(소제목) / item(①②③) / sub(※·- 보조)
STEPS = {
    1: ("사규 체계", [
        ("사규를 편(編) → 사규 → 별표·서식 3단계로 분류", "lead"),
        ("① 사규 목록 확인: 편(예: 2편 직제·윤리) 클릭", "item"),
        ("② 별표·서식 확인: [ + ] 버튼 클릭", "item"),
    ]),
    2: ("통합 검색창 (1)", [
        ("① 통합 검색창에 키워드 입력하여 내용 검색", "item"),
    ]),
    3: ("통합 검색창 (2)", [
        ("[검색어 입력 시 표시되는 결과 화면]", "lead"),
        ("① 결과 탭 : 사규 본문, 사규명, 별표, 부록 (탭명 옆 숫자 = 해당 탭 내 검색 결과 수)", "item"),
        ("② 결과 목록 : 검색 키워드와 결과 탭별 내 일치 건수 (N)를 표시 → 클릭 시 본문 조회 화면으로 이동", "item"),
    ]),
    4: ("최신 사규 제개정 및 폐지 목록", [
        ("① 사규 개정: 최근 개정된 사규 목록 확인", "item"),
        ("② 사규 제정·폐지: 최근 제정 또는 폐지된 사규 목록 확인", "item"),
        ("※ 사규명 클릭 시 사규 본문 조회 화면으로 이동합니다", "sub"),
        ("③ [ 더보기 ] 버튼 클릭 시 사규 캘린더 화면으로 이동", "item"),
    ]),
    5: ("공지사항 및 별표·부록 목록", [
        ("① 공지사항: 최근 등록된 공지사항 목록 확인", "item"),
        ("② 별표·부록: 최근 등록된 별표·부록 목록 확인", "item"),
        ("③ [ 더보기 ] 버튼 클릭 시 공지사항 및 별표·부록 전체 목록 화면으로 이동", "item"),
    ]),
    6: ("사규 목차", [
        ("① 장(章) → 절(節) → 조(條) 구조로 구성되고, 펼침/접힘(▼) 기능 활용 가능 ([모두 접기] 지원)", "item"),
        ("② 조문 클릭 시 본문 해당 위치로 이동", "item"),
    ]),
    7: ("사규 정보", [
        ("① 사규 기본 정보 확인 (소관부서명 / 제개정 시행일자 / 별표·부록 개수)", "item"),
        ("② [상세보기] 클릭 후 사규 상세 정보 확인 (제정/개정/시행일자, 문서번호, 제개정사유, 소관부서)", "item"),
    ]),
    8: ("기능 버튼", [
        ("우측 상단 기능 버튼 종류", "lead"),
        ("① [개정이력] : 과거 제정·개정·폐지 이력 조회", "item"),
        ("② [본문 검색] : 본문 내 키워드 검색", "item"),
        ("③ [가- / 원래대로 / 가+] : 글자 크기 조절", "item"),
        ("④ [2단보기] : 본문을 2단으로 보기", "item"),
        ("⑤ [신구대비] : 신·구 조문 대비표 확인", "item"),
        ("⑥ [전문다운] : 전체 본문 다운로드", "item"),
        ("⑦ [안내] : 사규 관련 안내사항 확인", "item"),
    ]),
    9: ("캘린더 기능", [
        ("① 해당 '연도' 선택", "item"),
        ("② 해당 '월' 선택", "item"),
        ("③ [오늘] 버튼 클릭 시, 당일 날짜로 자동 이동", "item"),
        ("④ 날짜 칸에 해당일자 시행 사규 표시", "item"),
    ]),
    10: ("사규 목록", [
        ("① 날짜 선택", "item"),
        ("② 해당 날짜에 제개정 시행된 사규가 목록에 표시", "item"),
        ("③ 사규명 클릭 시 사규 본문 조회 화면으로 이동", "item"),
    ]),
    11: ("사규 검색", [
        ("① 검색창에 키워드 입력 후 검색", "item"),
        ("② 날짜와 상관없이 검색한 사규 표시", "item"),
    ]),
    12: ("'법령 검색' 클릭", [
        ("시작 화면에서 [법령 검색] 클릭 시 법령 검색 화면으로 이동 (다음 페이지 참조)", "lead"),
    ]),
    13: ("법령 검색 화면", [
        ("① 검색 구분 선택 가능 (법령명 or 조문 내용)", "item"),
        ("② 키워드 검색", "item"),
        ("③ 최신 시행일자 기준 국가 법령 목록 확인 (키워드 입력 시 키워드 기준 목록 표시)", "item"),
    ]),
    14: ("법령 조회 화면", [
        ("① 검색 키워드 기준 법령 목록 표시", "item"),
        ("② 영역 클릭 시 본문 및 부칙 목차, 별표 목록 표시", "item"),
        ("③ 기능 버튼 : 조문 선택, 본문 검색, 본문 다운로드, 인쇄 등", "item"),
        ("④ 법령 정보 표시", "item"),
    ]),
    15: ("법령 본문 다운로드", [
        ("① 범위설정 (일부/선택 저장 · 조문 선택 · 포함 옵션)", "item"),
        ("- 전체 또는 조문 선택을 통한 다운로드 모두 가능", "sub"),
        ("- 출력 형식 선택 가능", "sub"),
        ("② 용지/폰트 설정 (본문 폰트 및 여백 설정 옵션)", "item"),
    ]),
}

STEP_PART = {}
for roman, ptitle, nums in PARTS:
    for n in nums:
        STEP_PART[n] = (roman, ptitle)


# ------------------------------------------------------------ 슬라이드 헬퍼
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_number_field(paragraph, size, color):
    """문단에 자동 슬라이드 번호(slidenum) 필드 추가 — 슬라이드 추가/삭제 시 자동 갱신."""
    p = paragraph._p
    fld = p.makeelement(qn('a:fld'), {'id': '{' + str(uuid.uuid4()).upper() + '}', 'type': 'slidenum'})
    rPr = p.makeelement(qn('a:rPr'), {'lang': 'ko-KR', 'sz': str(int(size * 100)), 'b': '1'})
    sf = p.makeelement(qn('a:solidFill'), {})
    sf.append(p.makeelement(qn('a:srgbClr'), {'val': str(color)}))
    rPr.append(sf)
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        rPr.append(p.makeelement(qn(tag), {'typeface': KR_FONT}))
    t = p.makeelement(qn('a:t'), {}); t.text = '0'
    fld.append(rPr); fld.append(t)
    p.append(fld)


def footer(slide):
    fy = SLIDE_H_IN - 0.52
    # 좌측: KB신용정보 로고(표지와 동일 마크)
    rect(slide, 0.45, fy + 0.05, 0.22, 0.22, fill=KB_YELLOW)
    add_text(slide, 0.74, fy, 6, 0.32,
             [[("KB신용정보 사규관리시스템", dict(size=12, bold=True, color=GREY_DK))]],
             anchor=MSO_ANCHOR.MIDDLE)
    # 우측: 자동 페이지 번호(12pt)
    tb = slide.shapes.add_textbox(Inches(SLIDE_W_IN - 1.35), Inches(fy), Inches(0.9), Inches(0.32))
    tb.text_frame.word_wrap = False
    pp = tb.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    _slide_number_field(pp, 12, GREY)


def marker(slide, cx, cy, d, n):
    """이미지와 분리된 번호 마커(별도 도형). cx,cy=중심(in), d=지름(in)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = MARK_RED
    sp.line.color.rgb = WHITE; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    style_run(r, size=9, bold=True, color=WHITE)


def desc_box(slide, x, y, w, h, items):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for text, kind in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = text
        if kind == "lead":
            style_run(r, size=12, bold=True, color=KB_YELLOW_DK); p.space_after = Pt(7)
        elif kind == "sub":
            style_run(r, size=10, color=GREY); p.space_after = Pt(5); p.level = 1
        else:
            style_run(r, size=11.5, color=GREY_DK); p.space_after = Pt(7)
    return tb


def cover_slide(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    rect(s, 0.7, 0.7, 0.34, 0.34, fill=KB_YELLOW)
    add_text(s, 1.15, 0.66, 6, 0.45, [[(f"{DOC['org']} 사규관리시스템", dict(size=15, bold=True, color=GREY_DK))]],
             anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 0, 2.35, SLIDE_W_IN, 2.35, fill=KB_YELLOW)
    rect(s, 0, 2.35, SLIDE_W_IN, 0.08, fill=KB_YELLOW_DK)
    add_text(s, 0.9, 2.72, 11.5, 1.0, [[(DOC["title"], dict(size=38, bold=True, color=WHITE))]])
    add_text(s, 0.9, 3.82, 11.5, 0.6, [[(f"Ver {DOC['version'].lstrip('vV')}", dict(size=20, bold=True, color=GREY_DK))]])
    # 하단 정보 박스
    bx, by, bw, bh = 0.9, 5.2, 7.6, 1.7
    rect(s, bx, by, bw, bh, fill=GREY_LT, line=GREY_BORDER, line_w=1)
    info = [
        [("작성 부서 : ", dict(size=11, bold=True, color=GREY_DK)), (DOC["dept"], dict(size=11, color=GREY_DK))],
        [("담당자 : ", dict(size=11, bold=True, color=GREY_DK)), (DOC["authors"], dict(size=11, color=GREY_DK))],
        [("작성 일자 : ", dict(size=11, bold=True, color=GREY_DK)), (DOC["written"], dict(size=11, color=GREY_DK)),
         ("    수정 일자 : ", dict(size=11, bold=True, color=GREY_DK)), (DOC["revised"], dict(size=11, color=GREY_DK))],
        [(f"* {DOC['contact']}", dict(size=9.5, color=GREY))],
    ]
    add_text(s, bx + 0.25, by + 0.2, bw - 0.5, bh - 0.4, info)
    add_text(s, SLIDE_W_IN - 3.0, 6.55, 2.6, 0.35,
             [[(f"{DOC['org']}  |  {DOC['steps']}단계", dict(size=10, color=GREY))]], align=PP_ALIGN.RIGHT)
    return s


def divider_slide(prs, roman, title):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    rect(s, 0, 2.75, SLIDE_W_IN, 2.0, fill=KB_YELLOW)
    rect(s, 0, 2.75, SLIDE_W_IN, 0.07, fill=KB_YELLOW_DK)
    add_text(s, 0.9, 3.05, 2.2, 1.4, [[(roman, dict(size=54, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 2.6, 3.05, 9.5, 1.4, [[(title, dict(size=34, bold=True, color=GREY_DK))]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s


def step_slide(prs, n, captures_dir, tmp):
    title, items = STEPS[n]
    roman, ptitle = STEP_PART[n]
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    # 헤더: [번호] 단계 제목(좌, 강조) ........... 섹션(우, 보조)
    rect(s, 0, 0, SLIDE_W_IN, 1.0, fill=KB_YELLOW)
    rect(s, 0, 1.0, SLIDE_W_IN, 0.05, fill=KB_YELLOW_DK)
    number_chip(s, 0.5, 0.27, 0.46, n, fill=WHITE, txt=KB_YELLOW_DK, size=13)
    add_text(s, 1.12, 0.1, 8.3, 0.8, [[(title, dict(size=22, bold=True, color=WHITE))]],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 9.5, 0.1, 3.45, 0.8, [[(f"{roman}. {ptitle}", dict(size=13, bold=True, color=WHITE))]],
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    # 좌측 캡처(확대): 거의 좌측 전체 폭
    IMGX, IMGY, IMGW, IMGH = 0.1, 1.45, 8.4, 5.5
    src = os.path.join(captures_dir, f"step{n:02d}.png") if captures_dir else None
    out_img = os.path.join(tmp, f"step{n:02d}.png")
    iw, ih = build_image(src, [], out_img)  # 깨끗한 원본(마커는 별도 도형)
    scale = min(IMGW / (iw / 96.0), IMGH / (ih / 96.0))
    dw, dh = (iw / 96.0) * scale, (ih / 96.0) * scale
    dx, dy = IMGX + (IMGW - dw) / 2, IMGY + (IMGH - dh) / 2
    has = bool(src and os.path.exists(src))
    if not has:
        rect(s, dx, dy, dw, dh, fill=GREY_LT, line=GREY_BORDER, line_w=1)
    s.shapes.add_picture(out_img, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    rect(s, dx, dy, dw, dh, fill=None, line=GREY_BORDER, line_w=1)  # 얇은 테두리
    if not has:
        add_text(s, dx, dy + dh / 2 - 0.35, dw, 0.7,
                 [[("[ 화면 캡처 영역 ]", dict(size=14, bold=True, color=GREY))],
                  [(f"step{n:02d} 캡처가 들어갈 자리", dict(size=10, color=GREY))]],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # 번호 마커(이미지와 분리된 별도 도형, 작게)
    for mn, fx, fy in MARKERS.get(n, []):
        marker(s, dx + fx * dw, dy + fy * dh, 0.23, mn)

    # 구분선 + 우측 설명(오른쪽으로 이동)
    rect(s, 8.62, 1.5, 0.014, 5.4, fill=GREY_BORDER)
    RX, RW = 8.85, 4.3
    add_text(s, RX, 1.5, RW, 0.35, [[("사용 방법", dict(size=13, bold=True, color=KB_YELLOW_DK))]])
    rect(s, RX, 1.87, 0.55, 0.045, fill=KB_YELLOW)
    desc_box(s, RX, 2.05, RW, 4.7, items)
    footer(s)
    return s


def build(captures_dir, out_path):
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * EMU_PER_IN))
    prs.slide_height = Emu(int(SLIDE_H_IN * EMU_PER_IN))
    cover_slide(prs)
    with tempfile.TemporaryDirectory() as tmp:
        for roman, ptitle, nums in PARTS:
            divider_slide(prs, roman, ptitle)
            for n in nums:
                step_slide(prs, n, captures_dir, tmp)
        prs.save(out_path)
    n_cap = 0
    if captures_dir and os.path.isdir(captures_dir):
        n_cap = sum(1 for n in STEPS if os.path.exists(os.path.join(captures_dir, f"step{n:02d}.png")))
    total = 1 + len(PARTS) + len(STEPS)
    print(f"생성 완료: {out_path}  (슬라이드 {total}장, 캡처 매칭 {n_cap}/{len(STEPS)})")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--captures", default="captures/stephow")
    ap.add_argument("--out", default="사용자매뉴얼/생성본/사규관리시스템_사용자매뉴얼_v0.1.pptx")
    args = ap.parse_args()
    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    cap = args.captures if os.path.isdir(args.captures) else None
    build(cap, args.out)
