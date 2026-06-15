# -*- coding: utf-8 -*-
"""사규관리시스템 사용자 매뉴얼 PPTX 생성기.

사용법:
    python tools/build_manual_pptx.py [--captures CAPTURES_DIR] [--out OUTPUT.pptx]

- CAPTURES_DIR 에 화면 캡처(png/jpg)를 두면 manual_content.py 의 image 파일명과
  매칭하여 슬라이드 좌측에 삽입합니다. 캡처가 없으면 자리표시자가 들어갑니다.
- 캡처 위에는 콜아웃 번호 마커가, 우측에는 같은 번호의 설명이 배치됩니다.

디자인: KB 톤(노랑+그레이), 16:9, 1화면=1슬라이드(좌 캡처/우 단계).
"""
import argparse
import os
import tempfile

from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import manual_content as MC

# ---------------------------------------------------------------- 테마/상수
KB_YELLOW    = RGBColor(0xFF, 0xB8, 0x00)
KB_YELLOW_DK = RGBColor(0xE8, 0x9A, 0x00)
GREY_DK      = RGBColor(0x33, 0x33, 0x33)
GREY         = RGBColor(0x6E, 0x6E, 0x6E)
GREY_LT      = RGBColor(0xF4, 0xF4, 0xF4)
GREY_BORDER  = RGBColor(0xDD, 0xDD, 0xDD)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
NOTE_RED     = RGBColor(0xD9, 0x53, 0x2C)

KR_FONT = "맑은 고딕"  # PowerPoint(Windows) 기본 한글 폰트

EMU_PER_IN = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# 주석 숫자용 TTF (Pillow 렌더)
_NUM_FONT_PATHS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
]


def _num_font(size):
    for p in _NUM_FONT_PATHS:
        if os.path.exists(p):
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


# ---------------------------------------------------------------- 폰트 헬퍼
def style_run(run, size, bold=False, color=GREY_DK, font=KR_FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    # 한글(East Asian) 글꼴 명시
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def add_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
             word_wrap=True):
    """lines: [(text, {size,bold,color,font}), ...] 또는 [[run,run], ...] 문단 리스트."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for para in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        runs = para if isinstance(para, list) else [para]
        for text, st in runs:
            r = p.add_run()
            r.text = text
            style_run(r, **st)
        if isinstance(para, list) and para and isinstance(para[-1][1], dict):
            sp = para[-1][1].get("space_after")
            if sp:
                p.space_after = Pt(sp)
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    return sp


def number_chip(slide, x, y, d, n, fill=KB_YELLOW, txt=GREY_DK, size=12):
    sp = rect(slide, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    tf = sp.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    style_run(r, size=size, bold=True, color=txt)
    return sp


# ---------------------------------------------------------------- 이미지(캡처+주석)
def build_image(src_path, callouts, out_path):
    """캡처(또는 자리표시자) 위에 콜아웃 번호 마커/박스를 그려 저장. (w,h) 반환."""
    if src_path and os.path.exists(src_path):
        img = Image.open(src_path).convert("RGB")
    else:
        img = _placeholder(1600, 940)
    W, H = img.size
    draw = ImageDraw.Draw(img, "RGBA")

    # 자동 배치용 좌표 (좌표 미지정 콜아웃)
    auto = [c for c in callouts if "x" not in c or "y" not in c]
    for i, c in enumerate(auto):
        c.setdefault("_ax", 0.09)
        c.setdefault("_ay", 0.16 + i * (0.66 / max(1, len(auto))))

    r = max(18, int(W * 0.020))
    fnt = _num_font(int(r * 1.25))
    for c in callouts:
        cx = int(c.get("x", c.get("_ax", 0.1)) * W)
        cy = int(c.get("y", c.get("_ay", 0.1)) * H)
        if "box" in c:
            bx, by, bw, bh = c["box"]
            draw.rectangle([bx * W, by * H, (bx + bw) * W, (by + bh) * H],
                           outline=(232, 154, 0, 255), width=max(3, int(W * 0.003)))
        # 마커: 흰 테두리 + 노랑 원 + 숫자
        draw.ellipse([cx - r - 3, cy - r - 3, cx + r + 3, cy + r + 3], fill=(255, 255, 255, 235))
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(255, 184, 0, 255))
        tb = draw.textbbox((0, 0), str(c["n"]), font=fnt)
        tw, th = tb[2] - tb[0], tb[3] - tb[1]
        draw.text((cx - tw / 2 - tb[0], cy - th / 2 - tb[1]), str(c["n"]), font=fnt, fill=(51, 51, 51, 255))
    img.save(out_path)
    return img.size


def _placeholder(W, H):
    img = Image.new("RGB", (W, H), (247, 247, 248))
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, W - 1, H - 1], outline=(214, 214, 214), width=4)
    # 가운데 옅은 안내 박스(테두리만) — 한글 텍스트는 PPTX 오버레이로 처리
    bw, bh = int(W * 0.42), int(H * 0.20)
    bx, by = (W - bw) // 2, (H - bh) // 2
    d.rounded_rectangle([bx, by, bx + bw, by + bh], radius=14, outline=(208, 208, 208), width=3)
    return img


# ---------------------------------------------------------------- 슬라이드
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def footer(slide, page):
    add_text(slide, 0.45, SLIDE_H_IN - 0.42, 8, 0.3,
             [[(MC.DOC["footer"], dict(size=8, color=GREY))]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W_IN - 1.2, SLIDE_H_IN - 0.42, 0.75, 0.3,
             [[(str(page), dict(size=8, color=GREY))]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def title_slide(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    # KB 마크 자리(노랑 사각) + 시스템명
    rect(s, 0.7, 0.7, 0.34, 0.34, fill=KB_YELLOW)
    add_text(s, 1.15, 0.66, 6, 0.45, [[(MC.DOC["system"], dict(size=15, bold=True, color=GREY_DK))]],
             anchor=MSO_ANCHOR.MIDDLE)
    # 중앙 노랑 밴드
    rect(s, 0, 2.55, SLIDE_W_IN, 2.5, fill=KB_YELLOW)
    rect(s, 0, 2.55, SLIDE_W_IN, 0.08, fill=KB_YELLOW_DK)
    add_text(s, 0.9, 2.95, 11.5, 1.0, [[(MC.DOC["system"], dict(size=40, bold=True, color=WHITE))]])
    add_text(s, 0.9, 3.95, 11.5, 0.8, [[(MC.DOC["title"], dict(size=30, bold=True, color=GREY_DK))]])
    # 하단 정보
    add_text(s, 0.9, 5.5, 11.5, 0.4, [[(f"접속 주소 : {MC.DOC['url']}", dict(size=13, color=GREY))]])
    add_text(s, 0.9, 5.95, 11.5, 0.4, [[(f"{MC.DOC['date']}", dict(size=13, color=GREY))]])
    add_text(s, 0.9, 6.6, 11.5, 0.4,
             [[("※ 본 매뉴얼은 사용자(임직원) 기준이며, 실제 화면과 일부 다를 수 있습니다.",
                dict(size=10, color=GREY))]])
    return s


def toc_slide(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    rect(s, 0, 0, SLIDE_W_IN, 1.1, fill=KB_YELLOW)
    rect(s, 0, 1.1, SLIDE_W_IN, 0.06, fill=KB_YELLOW_DK)
    add_text(s, 0.7, 0, 10, 1.1, [[("목차", dict(size=26, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
    # 2열 배치
    col_x = [0.8, 7.0]
    y0 = 1.7
    per = (len(MC.SCREENS) + 1) // 2
    for i, sc in enumerate(MC.SCREENS):
        col = i // per
        row = i % per
        x = col_x[col]
        y = y0 + row * 0.92
        number_chip(s, x, y, 0.46, sc["id"], size=12)
        tag = "" if sc.get("confirmed") else "  (확인 필요)"
        add_text(s, x + 0.65, y - 0.05, 5.2, 0.6,
                 [[(sc["name"], dict(size=14, bold=True, color=GREY_DK)),
                   (tag, dict(size=10, color=NOTE_RED))]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 2)
    return s


def screen_slide(prs, sc, page, captures_dir, tmpdir):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=WHITE)
    # 헤더
    rect(s, 0, 0, SLIDE_W_IN, 1.0, fill=KB_YELLOW)
    rect(s, 0, 1.0, SLIDE_W_IN, 0.05, fill=KB_YELLOW_DK)
    number_chip(s, 0.55, 0.27, 0.46, sc["id"], fill=WHITE, txt=KB_YELLOW_DK, size=13)
    add_text(s, 1.2, 0, 9, 1.0, [[(sc["name"], dict(size=22, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)

    # 좌측 캡처 영역
    LX, LY, LW, LH = 0.45, 1.35, 7.35, 5.55
    rect(s, LX, LY, LW, LH, fill=GREY_LT, line=GREY_BORDER, line_w=1)
    src = os.path.join(captures_dir, sc["image"]) if captures_dir else None
    out_img = os.path.join(tmpdir, f"{sc['id']}.png")
    iw, ih = build_image(src, sc.get("callouts", []), out_img)
    # 비율 맞춰 박스 안에 배치
    pad = 0.12
    bw, bh = LW - 2 * pad, LH - 2 * pad
    scale = min(bw / (iw / 96.0), bh / (ih / 96.0))
    dw, dh = (iw / 96.0) * scale, (ih / 96.0) * scale
    dx, dy = LX + (LW - dw) / 2, LY + (LH - dh) / 2
    s.shapes.add_picture(out_img, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    has_capture = bool(src and os.path.exists(src))
    if not has_capture:
        add_text(s, LX, LY + LH / 2 - 0.35, LW, 0.7,
                 [[("[ 화면 캡처 영역 ]", dict(size=14, bold=True, color=GREY))],
                  [("(여기에 실제 캡처가 들어갑니다)", dict(size=10, color=GREY))]],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # 우측 설명 영역
    RX, RW = 8.05, 4.85
    add_text(s, RX, 1.45, RW, 0.35, [[("화면 개요", dict(size=13, bold=True, color=KB_YELLOW_DK))]])
    rect(s, RX, 1.82, 0.55, 0.045, fill=KB_YELLOW)
    add_text(s, RX, 1.95, RW, 1.1, [[(sc["overview"], dict(size=11, color=GREY_DK))]])

    add_text(s, RX, 3.05, RW, 0.35, [[("사용 단계", dict(size=13, bold=True, color=KB_YELLOW_DK))]])
    rect(s, RX, 3.42, 0.55, 0.045, fill=KB_YELLOW)

    cy = 3.6
    for c in sc.get("callouts", []):
        number_chip(s, RX, cy, 0.32, c["n"], size=11)
        add_text(s, RX + 0.45, cy - 0.06, RW - 0.45, 0.7,
                 [[(c["text"], dict(size=10.5, color=GREY_DK))]], anchor=MSO_ANCHOR.TOP)
        # 줄 높이 추정(텍스트 길이 기반)
        cy += 0.34 + 0.20 * (len(c["text"]) // 26)

    if sc.get("note"):
        ny = max(cy + 0.05, 6.55)
        rect(s, RX, ny, RW, 0.42, fill=RGBColor(0xFD, 0xF1, 0xE6), line=NOTE_RED, line_w=0.75)
        add_text(s, RX + 0.12, ny, RW - 0.24, 0.42,
                 [[("[확인 필요] ", dict(size=9.5, bold=True, color=NOTE_RED)),
                   (sc["note"], dict(size=9.5, color=NOTE_RED))]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page)
    return s


def build(captures_dir, out_path):
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * EMU_PER_IN))
    prs.slide_height = Emu(int(SLIDE_H_IN * EMU_PER_IN))
    title_slide(prs)
    toc_slide(prs)
    with tempfile.TemporaryDirectory() as tmp:
        for i, sc in enumerate(MC.SCREENS):
            screen_slide(prs, sc, i + 3, captures_dir, tmp)
        prs.save(out_path)
    n_slides = 2 + len(MC.SCREENS)
    n_cap = 0
    if captures_dir and os.path.isdir(captures_dir):
        n_cap = sum(1 for sc in MC.SCREENS if os.path.exists(os.path.join(captures_dir, sc["image"])))
    print(f"생성 완료: {out_path}  (슬라이드 {n_slides}장, 캡처 매칭 {n_cap}/{len(MC.SCREENS)})")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--captures", default="captures", help="캡처 폴더 (기본: captures)")
    ap.add_argument("--out", default="사용자매뉴얼/생성본/KB사규관리시스템_사용자매뉴얼.pptx")
    args = ap.parse_args()
    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    cap = args.captures if os.path.isdir(args.captures) else None
    build(cap, args.out)
